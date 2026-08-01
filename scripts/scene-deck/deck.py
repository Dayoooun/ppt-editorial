# -*- coding: utf-8 -*-
"""씬 덱 단일 진입점 (2026-08-01)

## 왜 필요한가
프리셋·레이아웃·폰트·사진·수정 모듈이 다 갖춰졌는데, 정작 덱을 만들 때마다
`build.py` 보일러플레이트(스타일 조립 → jobs.json → 생성 → 레이아웃 세팅 → PDF)를
새로 써야 했다. 실측으로 6개 프로젝트가 같은 코드를 반복했고, 그때마다 편차가 생겼다.

## 사용법
```python
from deck import Deck

d = Deck(domain="식음료", foot="OO식품", title="사업 소개")
d.slide("L", "THE PROBLEM", ["산지는 좋은데", "팔 길이 없습니다"],
        ["좋은 원물을 확보해도 판로가 없으면 재고가 된다"],
        scene="a lone food stall on an isolated platform")
d.slide("W", "OUR PROCESS", ["산지에서 식탁까지"], ["4단계 일괄 관리"],
        scene="four pedestals in a row: crate, machine, package, delivery box",
        labels=["수급", "가공", "포장", "유통"])
d.slide("L", "TRACK RECORD", ["숫자로 증명합니다"], ["공공기관 검증"],
        scene="a tall stack of layered plates with check badges",
        num=["2,800", "회", "연간 커밋"], chips=["HACCP", "직거래"])
d.photos(["a.jpg", "b.jpg", "c.jpg"], "ON SITE", ["현장에서 함께합니다"], ["3회차 밀착"])

d.generate()          # 씬 생성 (없는 것만)
d.build()             # 조립 + PDF
```

한 번 쓴 덱은 `spec.json`으로 저장되어 `revise.py`로 수정할 수 있다.
"""
import os, sys, json, subprocess, glob

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

from presets import preset, style_block          # noqa: E402
import layout_engine as LE                        # noqa: E402

GEN = os.path.join(os.path.dirname(HERE), "codex_parallel_gen.py")

SAFE = ("\n★★ SAFE ZONES — OVERRIDES EVERYTHING:\n"
        "1) TOP 15% COMPLETELY EMPTY.  2) BOTTOM 15% COMPLETELY EMPTY.\n"
        "3) All content strictly between 15% and 85% of slide height.\n")

TAIL = ("\nUse ONLY the image-generation capability - ABSOLUTELY NO Python/PIL/code drawing.\n"
        "Korean text inside the image must render perfectly — no '?', no broken glyphs.\n"
        "결과를 반드시 __out.png 로 저장하고 크기를 출력하라. PIL 드로잉 금지, 이미지 생성만.\n")

# 구도별 권장 씬 비율
RATIO = {"W": "16:9 wide", "F": "16:9 wide"}
DEFAULT_RATIO = "4:3 landscape"


class Deck:
    def __init__(self, domain="it", foot="", title="", out_dir=None):
        self.domain = domain
        self.preset = preset(domain)
        self.foot = foot
        self.title = title
        self.dir = os.path.abspath(out_dir or "deck_out")
        self.slides = []
        self._photo_refs = {}
        os.makedirs(self.dir, exist_ok=True)

    # ══════ 슬라이드 추가 ══════
    def slide(self, lay, eyebrow, head, sub, scene=None, labels=None,
              num=None, chips=None, items=None, ratio=None):
        """슬라이드 1장 추가.

        lay      L/S/W/C/A/F/T
        scene    씬 내용(영문 서술). None이면 텍스트만.
        labels   씬 안에 박을 한글 라벨
        num      ["2,800", "회", "연간 커밋"] 대형 수치
        chips    ["키워드1", "키워드2"] 칩
        items    [("제목","설명"), ...] T 구도 전용
        """
        sid = "s%02d" % (len(self.slides) + 1)
        s = {"_sid": len(self.slides) + 1, "lay": lay, "eyebrow": eyebrow,
             "scene": sid, "head": list(head), "sub": list(sub)}
        if scene:
            s["scene_body"] = scene
            s["scene_ratio"] = ratio or RATIO.get(lay, DEFAULT_RATIO)
        if labels:
            s["scene_labels"] = list(labels)
        for k, v in [("num", num), ("chips", chips), ("items", items)]:
            if v:
                s[k] = v
        self.slides.append(s)
        self.save()          # 매 추가마다 저장 — generate 후 load 가능하게(실측 버그)
        return self


    # ══════ 실무 필수 슬라이드 ══════
    def cover(self, eyebrow, head, sub=None, issuer=None, meta=None, scene=None):
        """표지. 크롬(쪽번호·푸터)이 붙지 않는다."""
        self.slide("COVER", eyebrow, head, sub or [], scene=scene,
                   ratio="4:3 landscape")
        self.slides[-1].update({"issuer": issuer or self.foot,
                                "meta": meta or [], "_nochrome": True})
        self.save()
        return self

    def agenda(self, items, head=None, eyebrow="AGENDA", scene=None):
        """목차. items=["배경","솔루션",...] 또는 [("01","배경"),...]"""
        norm = [(("%02d" % i), x) if isinstance(x, str) else tuple(x)
                for i, x in enumerate(items, 1)]
        self.slide("AGENDA", eyebrow, head or ["목차"], [], scene=scene)
        self.slides[-1]["items"] = norm
        self.save()
        return self

    def closing(self, head, sub=None, issuer=None, scene=None, eyebrow="THANK YOU"):
        """클로징."""
        self.slide("CLOSING", eyebrow, head, sub or [], scene=scene)
        self.slides[-1]["issuer"] = issuer or self.foot
        self.save()
        return self

    def photos(self, paths, eyebrow, head, sub, lay=None, labels=None, **kw):
        """실사진 슬라이드. 철칙 D — refs로 codex에 전달해 씬 안에서 함께 그린다."""
        from photos import plan, prep, prompt_block
        ready = prep(paths, os.path.join(self.dir, "photos_ready"))
        pl = plan(ready)
        lay = lay or {"hero": "L", "compare": "C", "sequence": "W", "grid": "F"}[pl["mode"]]
        self.slide(lay, eyebrow, head, sub,
                   scene="a %s composition integrating the provided photographs" % pl["mode"],
                   labels=labels, **kw)
        sid = self.slides[-1]["scene"]
        self._photo_refs[sid] = (ready, pl, labels)
        self.slides[-1]["_photo"] = [os.path.relpath(x, self.dir) for x in ready]
        self.save()
        return self

    # ══════ 씬 생성 ══════
    def jobs(self, only_missing=True):
        """씬 생성 잡 목록. 이미 있는 씬은 건너뛴다."""
        from photos import prompt_block
        scn = os.path.join(self.dir, "scenes")
        os.makedirs(scn, exist_ok=True)
        out = []
        for s in self.slides:
            if not s.get("scene_body"):
                continue
            sid = s["scene"]
            path = os.path.join(scn, "%s.png" % sid)
            if only_missing and os.path.exists(path) and os.path.getsize(path) > 100_000:
                continue
            p = style_block(self.domain) + "\nCOMPOSITION: " + s["scene_body"]
            if sid in self._photo_refs:
                ready, pl, labels = self._photo_refs[sid]
                p += prompt_block(pl, labels)
                refs = [os.path.relpath(x, self.dir) for x in ready]
            else:
                refs = []
                if s.get("scene_labels"):
                    p += ("\n\nTEXT LABELS rendered in the image "
                          "(bold Korean gothic, ink colour, small):\n  "
                          + "  ".join('"%s"' % x for x in s["scene_labels"]))
            p += SAFE + "\nOutput %s framing.\n" % s.get("scene_ratio", DEFAULT_RATIO) + TAIL
            out.append({"label": sid, "refs": refs,
                        "out": os.path.join("scenes", "%s.png" % sid), "prompt": p})
        return out

    def generate(self, effort="high", timeout=800, loop=2):
        """씬 생성 실행. 이미 있는 씬은 재생성하지 않는다(철칙 E)."""
        jb = self.jobs()
        if not jb:
            print("[생성] 모든 씬이 이미 존재 — 생략")
            return self
        jp = os.path.join(self.dir, "jobs.json")
        json.dump(jb, open(jp, "w", encoding="utf-8"), ensure_ascii=False, indent=1)
        print("[생성] %d장 (effort=%s)" % (len(jb), effort))
        subprocess.run([sys.executable, GEN, "jobs.json", "--retry", "1",
                        "--loop", str(loop), "--timeout", str(timeout),
                        "--effort", effort], cwd=self.dir)
        # 생성 후 실제 산출 확인 — 실패한 씬을 조용히 넘기지 않는다
        failed = [j["label"] for j in jb
                  if not os.path.exists(os.path.join(self.dir, j["out"]))
                  or os.path.getsize(os.path.join(self.dir, j["out"])) < 100_000]
        if failed:
            print("  ⚠ 생성 실패 %d장: %s — 다시 generate() 하면 실패분만 재시도한다"
                  % (len(failed), ", ".join(failed)))
        self.save()
        return self

    # ══════ 조립 ══════
    def _apply_style(self):
        p = self.preset["palette"]
        LE.FAMILY = self.preset["fonts"]["head"]
        LE.SCN = os.path.join(self.dir, "scenes")
        LE.BLUE, LE.INK, LE.GREY, LE.LINE = p["hero"], p["ink"], p["grey"], p["line"]
        LE.FOOT = self.foot

    def build(self, pdf=True, pptx=False):
        """조립 + PDF. 씬이 없으면 텍스트만 렌더된다."""
        from PIL import Image, ImageDraw
        self._apply_style()
        out = os.path.join(self.dir, "out")
        os.makedirs(out, exist_ok=True)
        n = len(self.slides)
        for i, s in enumerate(self.slides, 1):
            im = Image.new("RGB", (LE.W, LE.H), LE.WHITE)
            d = ImageDraw.Draw(im)
            LE.LAY[s["lay"]](im, d, s)
            if not s.get("_nochrome"):
                LE.chrome(im, d, s["eyebrow"], "%02d" % i, n)
            im.save(os.path.join(out, "slide_%02d.png" % i))
        # 씬 누락 경고 — 조용히 텍스트만 렌더되면 납품물에 빈 슬라이드가 섞인다
        missing = [("%02d" % i, s.get("scene"))
                   for i, s in enumerate(self.slides, 1)
                   if s.get("scene") and not os.path.exists(
                       os.path.join(self.dir, "scenes", "%s.png" % s.get("scene")))]
        print("[조립] %d장 (도메인=%s 폰트=%s)" % (n, self.preset["key"], LE.FAMILY))
        if missing:
            print("  ⚠ 씬 누락 %d장 — 텍스트만 렌더됨: %s"
                  % (len(missing), ", ".join("%s(%s)" % m for m in missing)))
            print("    → d.generate() 로 재생성하라")
        self.save()
        outs = []
        if pdf:
            outs.append(self.pdf())
        if pptx:
            outs.append(self.pptx())
        return outs[0] if len(outs) == 1 else (outs or None)

    def pptx(self, name=None):
        """편집 가능한 PPTX 출력. 고객이 PPT 파일을 원할 때.
        슬라이드 전면에 PNG를 깐 형태라 텍스트 편집은 불가하지만
        발표 도구(슬라이드쇼·노트·인쇄)를 그대로 쓸 수 있다."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            print("[PPTX] python-pptx 미설치 — pip install python-pptx")
            return None
        out = os.path.join(self.dir, "out")
        fs = sorted(glob.glob(os.path.join(out, "slide_*.png")))
        if not fs:
            print("[PPTX] 슬라이드 없음 — build() 먼저")
            return None
        prs = Presentation()
        prs.slide_width = Inches(13.333)      # 16:9
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        for f in fs:
            sl = prs.slides.add_slide(blank)
            sl.shapes.add_picture(f, 0, 0, width=prs.slide_width,
                                  height=prs.slide_height)
        p = os.path.join(self.dir, (name or self.title or "deck") + ".pptx")
        prs.save(p)
        print("[PPTX] %s (%.1f MB)" % (p, os.path.getsize(p) / 1024 / 1024))
        return p

    def pdf(self, name=None):
        import fitz
        out = os.path.join(self.dir, "out")
        fs = sorted(glob.glob(os.path.join(out, "slide_*.png")))
        doc = fitz.open()
        for f in fs:
            pg = doc.new_page(width=1280, height=720)
            pg.insert_image(pg.rect, filename=f)
        p = os.path.join(self.dir, (name or self.title or "deck") + ".pdf")
        doc.save(p, deflate=True)
        print("[PDF] %s (%.1f MB)" % (p, os.path.getsize(p) / 1024 / 1024))
        return p

    # ══════ 저장/로드 (revise.py 연동) ══════
    def save(self):
        spec = {"domain": self.domain, "foot": self.foot, "title": self.title,
                "family": self.preset["fonts"]["head"], "slides": self.slides}
        p = os.path.join(self.dir, "spec.json")
        json.dump(spec, open(p, "w", encoding="utf-8"), ensure_ascii=False, indent=1)
        return p

    @classmethod
    def load(cls, spec_path):
        spec = json.load(open(spec_path, encoding="utf-8"))
        d = cls(domain=spec.get("domain", "it"), foot=spec.get("foot", ""),
                title=spec.get("title", ""), out_dir=os.path.dirname(spec_path))
        d.slides = spec["slides"]
        return d

    def revise(self):
        """revise.Deck 핸들 반환 — 한 줄 명령으로 수정"""
        from revise import Deck as RDeck
        return RDeck({"domain": self.domain, "foot": self.foot,
                      "family": self.preset["fonts"]["head"], "slides": self.slides},
                     os.path.join(self.dir, "spec.json"))

    def summary(self):
        print("=== %s (%s) ===" % (self.title or "덱", self.preset["key"]))
        for i, s in enumerate(self.slides, 1):
            ex = [k for k in ("num", "chips", "items") if s.get(k)]
            print("  %02d [%s] %-30s %s" % (i, s["lay"], s["head"][0][:30],
                                            "+".join(ex) or ""))


if __name__ == "__main__":
    print(__doc__.split("## 사용법")[1].strip()[:600])
