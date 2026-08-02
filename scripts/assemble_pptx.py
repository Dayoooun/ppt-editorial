"""ppt-hybrid — build_pptx.py
렌더된 PNG 폴더 → 16:9 PPTX (+ --pdf 시 PDF도).

사용:
  python build_pptx.py <png_dir> <out.pptx> [--pdf]
"""
import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("png_dir")
    ap.add_argument("out")
    ap.add_argument("--pdf", action="store_true")
    args = ap.parse_args()

    png_dir = Path(args.png_dir).resolve()
    pngs = sorted(png_dir.glob("*.png"))
    if not pngs:
        raise SystemExit(f"no png in {png_dir}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, prs.slide_width, prs.slide_height)
    prs.save(args.out)
    print(f"{len(pngs)} slide(s) -> {args.out}")

    if args.pdf:
        import fitz
        doc = fitz.open()
        for png in pngs:
            img = fitz.open(str(png))
            rect = img[0].rect
            pdf_bytes = img.convert_to_pdf()
            img.close()
            src = fitz.open("pdf", pdf_bytes)
            page = doc.new_page(width=rect.width, height=rect.height)
            page.show_pdf_page(page.rect, src, 0)
        pdf_out = str(Path(args.out).with_suffix(".pdf"))
        # [2026-07-29] 뷰어가 PDF를 열어 두면 저장이 Permission denied로 죽는다.
        # 스택 트레이스 대신 무엇을 해야 하는지 한 줄로 말한다.
        try:
            doc.save(pdf_out)
        except Exception as e:
            if "Permission denied" in str(e) or "denied" in str(e).lower():
                sys.exit(f"PDF를 저장할 수 없다 — {pdf_out} 가 뷰어에 열려 있다. 닫고 다시 실행해라")
            raise
        print(f"pdf -> {pdf_out}")


if __name__ == "__main__":
    main()
